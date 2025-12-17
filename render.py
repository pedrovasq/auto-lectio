from __future__ import annotations

import argparse
import os
from datetime import datetime
import json
from copy import deepcopy
from pathlib import Path
from typing import Dict, List, Tuple

from pptx import Presentation


# Simple iterator that also descends into group shapes (type 6)
def iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        # 6 == MSO_SHAPE_TYPE.GROUP
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            for s in shape.shapes:
                yield s


def _replace_in_text_frame(tf, token: str, new_text: str) -> bool:
    replaced = False
    for p in tf.paragraphs:
        for run in p.runs:
            if token in run.text:
                run.text = run.text.replace(token, new_text)
                replaced = True
    return replaced


def replace_token_in_shape(shape, token: str, new_text: str) -> bool:
    """Replace token in a single shape (text frame or table cells)."""
    replaced = False
    if getattr(shape, "has_text_frame", False):
        if _replace_in_text_frame(shape.text_frame, token, new_text):
            replaced = True
    # Tables
    if getattr(shape, "has_table", False):
        tbl = shape.table
        for row in tbl.rows:
            for cell in row.cells:
                if _replace_in_text_frame(cell.text_frame, token, new_text):
                    replaced = True
    return replaced


def replace_tokens_in_slide(slide, mapping: Dict[str, str]) -> int:
    """Replace all tokens from mapping in the given slide. Returns count replaced."""
    count = 0
    for shape in iter_shapes(slide):
        for token, val in mapping.items():
            if replace_token_in_shape(shape, token, val):
                count += 1
    return count


def slide_contains_token(slide, token: str) -> bool:
    for shape in iter_shapes(slide):
        if getattr(shape, "has_text_frame", False) and token in shape.text_frame.text:
            return True
        if getattr(shape, "has_table", False):
            tbl = shape.table
            for row in tbl.rows:
                for cell in row.cells:
                    if token in cell.text_frame.text:
                        return True
    return False


def _sanitize_text(s: str) -> str:
    if s is None:
        return ""
    # Replace all newlines with spaces and collapse repeated whitespace
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    s = s.replace("\n", " ")
    # Collapse multiple spaces/tabs
    s = " ".join(s.split())
    return s.strip()


def enforce_chunk_bounds(chunks: List[str], min_chars: int = 100, max_chars: int = 140) -> List[str]:
    """Merge adjacent short chunks when possible to keep within [min,max].

    - Only merges with the immediate next chunk if combined length <= max.
    - Leaves chunks as-is when merging would exceed max.
    - Operates on already-sanitized chunks (no newlines).
    """
    out: List[str] = []
    i = 0
    n = len(chunks)
    while i < n:
        cur = chunks[i]
        # If current is short and there is a next, try to merge
        if len(cur) < min_chars and i + 1 < n:
            nxt = chunks[i + 1]
            # Prefer to merge if it keeps us within max
            if len(cur) + 1 + len(nxt) <= max_chars:
                cur = (cur + " " + nxt).strip()
                i += 2
                out.append(cur)
                continue
        out.append(cur)
        i += 1
    return out


def find_seed_slide_indices(prs: Presentation, token: str) -> List[int]:
    return [i for i, s in enumerate(prs.slides) if slide_contains_token(s, token)]


def insert_slide_after(prs: Presentation, new_slide, insert_after_index: int) -> None:
    """Move the specific slide (by part rel) to position right after insert_after_index."""
    sldIdLst = prs.slides._sldIdLst
    # Find relationship id for this slide
    target_part = new_slide.part
    rId = None
    for rel in prs.part.rels.values():
        if getattr(rel, "_target", None) is target_part:
            rId = rel.rId
            break
    if rId is None:
        # Fallback to last element (best-effort)
        new_id = sldIdLst[-1]
        sldIdLst.remove(new_id)
        sldIdLst.insert(insert_after_index + 1, new_id)
        return
    # Find the sldId element with that rId
    moving = None
    for sldId in sldIdLst:
        if sldId.rId == rId:
            moving = sldId
            break
    if moving is None:
        return
    sldIdLst.remove(moving)
    sldIdLst.insert(insert_after_index + 1, moving)


def _shape_text(shape) -> str:
    try:
        if getattr(shape, "has_text_frame", False):
            return shape.text_frame.text or ""
        if getattr(shape, "has_table", False):
            parts = []
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text_frame.text or "")
            return "\n".join(parts)
    except Exception:
        return ""
    return ""

def clear_shapes(slide) -> None:
    for shape in list(slide.shapes):
        el = shape._element
        el.getparent().remove(el)


def _shape_has_image_rel(shape) -> bool:
    try:
        # detect any blip reference which would need a rel copy
        blips = shape._element.xpath('.//a:blip', namespaces=shape._element.nsmap)
        return bool(blips)
    except Exception:
        return False


def duplicate_slide_filtered(prs: Presentation, seed_index: int, current_key: str, known_tokens: set[str]) -> Tuple[int, object]:
    """Duplicate a slide by copying only safe shapes to avoid repair prompts and stray placeholders.

    - Copies text shapes containing the current_key token.
    - Copies shapes with no placeholder tokens in their text.
    - Skips shapes that contain any other placeholder tokens.
    - Skips shapes with image relationships to avoid missing rels/repair.
    """
    seed = prs.slides[seed_index]
    new_slide = prs.slides.add_slide(seed.slide_layout)
    insert_slide_after(prs, new_slide, seed_index)
    clear_shapes(new_slide)

    def should_copy(shape) -> bool:
        if _shape_has_image_rel(shape):
            return False
        txt = _shape_text(shape)
        if current_key in txt:
            return True
        if any(tok in txt for tok in known_tokens if tok != current_key):
            return False
        # copy everything else (title boxes, static labels, shapes without text)
        return True

    for shape in seed.shapes:
        if should_copy(shape):
            new_el = deepcopy(shape._element)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return seed_index + 1, new_slide


def load_payload(json_path: str) -> Dict:
    with open(json_path, "r", encoding="utf-8") as f:
        return json.load(f)


def chunk_psalm_text(text: str) -> List[str]:
    """Split psalm into alternating refrain and verse blocks.

    Pattern: Each line starting with 'R.' (optionally with qualifiers like '(7a)')
    is its own chunk; the following non-R lines up to the next R form a verse chunk.
    """
    if text is None:
        return []
    lines = [ln.strip() for ln in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    lines = [ln for ln in lines if ln]
    chunks: List[str] = []
    current_verse: List[str] = []
    import re
    is_refrain = lambda s: re.match(r"^R[\./]?(?:\s*\([^)]*\))?\s", s) is not None or s.startswith("R.")

    for ln in lines:
        if is_refrain(ln):
            # Flush verse collected so far
            if current_verse:
                verse = "\n".join(current_verse).strip()
                if verse:
                    chunks.append(verse)
                current_verse = []
            # Add refrain as its own chunk
            chunks.append(ln)
        else:
            current_verse.append(ln)

    if current_verse:
        verse = "\n".join(current_verse).strip()
        if verse:
            chunks.append(verse)

    # Ensure we start with refrain if present; otherwise leave as-is
    # Filter empties just in case
    return [c for c in chunks if c and c.strip()]


def delete_slides(prs: Presentation, indices: List[int]) -> None:
    """Delete slides by indices and drop their package relationships to avoid repair."""
    if not indices:
        return
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(set(indices), reverse=True):
        if 0 <= idx < len(sldIdLst):
            sldId = sldIdLst[idx]
            rId = sldId.rId
            prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)


def main() -> None:
    parser = argparse.ArgumentParser(description="Render PPTX from JSON payload with waterfall duplication.")
    parser.add_argument("--template", required=True, help="Path to template PPTX (e.g., template.pptx)")
    parser.add_argument("--json", dest="json_path", required=True, help="Path to payload JSON")
    parser.add_argument("--out", required=True, help="Path to output PPTX")
    parser.add_argument("--stamp", action="store_true", help="Append timestamp to output filename")
    parser.add_argument("--verbose", action="store_true", help="Enable verbose logging")
    args = parser.parse_args()

    payload = load_payload(args.json_path)
    placeholders: Dict[str, str] = payload.get("placeholders", {})
    chunks_map: Dict[str, List[str]] = payload.get("chunks", {})

    prs = Presentation(args.template)

    def log(msg: str) -> None:
        if args.verbose:
            print(msg)

    # Log initial positions of all placeholders before any replacement/deletion
    interested = [
        "{FIRST_READING_REF}", "{FIRST_READING_TXT}",
        "{PSALM_REF}", "{PSALM_TXT}",
        "{SECOND_READING_REF}", "{SECOND_READING_TXT}",
        "{ACCLAMATION_REF}", "{ACCLAMATION_TXT}",
        "{GOSPEL_REF}", "{GOSPEL_TXT}",
    ]
    if args.verbose:
        for tok in interested:
            idxs = find_seed_slide_indices(prs, tok)
            if idxs:
                log(f"Initial positions {tok}: {[i+1 for i in idxs]} (1-based)")

    # Define which placeholders should use waterfall expansion
    waterfall_keys = [
        "{FIRST_READING_TXT}",
        "{PSALM_TXT}",
        "{SECOND_READING_TXT}",
        "{GOSPEL_TXT}",
    ]

    # Known tokens for cleanup when missing
    known_tokens = {
        "{LITURGICAL_DAY}",
        "{FIRST_READING_REF}", "{FIRST_READING_TXT}",
        "{PSALM_REF}", "{PSALM_TXT}",
        "{SECOND_READING_REF}", "{SECOND_READING_TXT}",
        "{ACCLAMATION_REF}", "{ACCLAMATION_TXT}",
        "{GOSPEL_REF}", "{GOSPEL_TXT}",
    }

    # Simple replacements: everything except the waterfall keys
    simple_mapping = {k: _sanitize_text(v) for k, v in placeholders.items() if k not in waterfall_keys}
    # Cleanup: ensure missing non-waterfall tokens are blanked out
    for tok in known_tokens:
        if tok not in placeholders and tok not in waterfall_keys:
            simple_mapping[tok] = ""

    # For diagnostics: a superset including hymn placeholders
    log_tokens = set(interested + [
        "{ENTRANCE_HYMN}", "{OFFERTORY_HYMN}", "{MYSTERY_OF_FAITH}", "{COMMUNION_HYMN}", "{RECESSIONAL_HYMN}",
    ])

    def tokens_in_slide(slide) -> List[str]:
        present = []
        for tok in log_tokens:
            if slide_contains_token(slide, tok):
                present.append(tok)
        return sorted(present)

    def slide_text_preview(slide, limit: int = 120) -> str:
        parts: List[str] = []
        for shape in iter_shapes(slide):
            if getattr(shape, 'has_text_frame', False):
                t = shape.text_frame.text or ''
                if t:
                    parts.append(t.replace('\n', '|'))
            if getattr(shape, 'has_table', False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        ct = cell.text_frame.text or ''
                        if ct:
                            parts.append(ct.replace('\n', '|'))
        txt = ' || '.join(parts)
        if len(txt) > limit:
            return txt[:limit] + '...'
        return txt

    def snapshot(label: str) -> None:
        if not args.verbose:
            return
        print(f"SNAP[{label}] total_slides={len(prs.slides)}")
        for idx, slide in enumerate(prs.slides, start=1):
            toks = tokens_in_slide(slide)
            if toks:
                print(f"  slide {idx}: tokens={toks}")
            else:
                # still show a short text preview to inspect blanks
                prev = slide_text_preview(slide)
                print(f"  slide {idx}: tokens=[], preview='{prev}'")

    # Snapshot before any changes
    snapshot('before')

    # If there is no second reading, do NOT delete slides (can corrupt package);
    # instead, blank placeholders via simple_mapping cleanup below.
    sec_txt = (placeholders.get("{SECOND_READING_TXT}") or "").strip()
    sec_chunks = chunks_map.get("{SECOND_READING_TXT}") or []
    has_second_reading = bool(sec_txt) or (isinstance(sec_chunks, list) and any((c or "").strip() for c in sec_chunks))
    # log status but keep slides
    if not has_second_reading:
        log("No second reading detected; leaving slides in place and blanking placeholders.")

    # 1) Replace simple placeholders across all slides
    for slide in prs.slides:
        replace_tokens_in_slide(slide, simple_mapping)

    # 2) Waterfall expansion for long body text
    # Find seed indices first
    seeds: List[Tuple[str, int]] = []
    for key in waterfall_keys:
        indices = find_seed_slide_indices(prs, key)
        if not indices:
            # No seed present for this key; fall back to replacing token as-is
            val = _sanitize_text(placeholders.get(key, ""))
            log(f"No seed for {key}; applying simple replacement across deck")
            for slide in prs.slides:
                replace_tokens_in_slide(slide, {key: val})
            continue
        # assume exactly one seed per key; use the first if multiple
        seed_idx = indices[0]
        seeds.append((key, seed_idx))
        log(f"Seed for {key} at slide index {seed_idx}")

    # Process in descending order to avoid shifting indices of future seeds
    seeds.sort(key=lambda kv: kv[1], reverse=True)

    # Process each seed in order
    for key, seed_index in seeds:
        # Prepare chunks (fallback to raw placeholder text if chunks missing)
        if key == "{PSALM_TXT}":
            # Always derive psalm chunks from the full text to ensure correct alternation
            raw = placeholders.get(key, "") or ""
            chunks = chunk_psalm_text(raw)
        else:
            chunks = chunks_map.get(key)
            if not chunks:
                raw = placeholders.get(key, "")
                chunks = [raw] if raw is not None else [""]

        # Filter out empty chunks to avoid blank slides and sanitize whitespace
        chunks = [_sanitize_text(c) for c in chunks if c and c.strip()]
        # Enforce desired bounds for non-psalm waterfalls
        if key != "{PSALM_TXT}":
            chunks = enforce_chunk_bounds(chunks, min_chars=100, max_chars=140)
        log(f"{key}: {len(chunks)} chunk(s)")

        if len(chunks) == 0:
            continue

        # Duplicate slides for all chunks beyond the first
        # Strategy: create N-1 new slides right after the current tail of this sequence
        log(f"{key}: seed slide {seed_index+1} tokens before dup: {tokens_in_slide(prs.slides[seed_index])}")
        current_index = seed_index
        created_indices: List[int] = []
        # Create N-1 duplicates of the seed, inserted sequentially after it
        for _ in range(len(chunks) - 1):
            new_index, new_slide = duplicate_slide_filtered(prs, current_index, key, known_tokens)
            # Fill simple placeholders on the newly created slide as well
            replace_tokens_in_slide(new_slide, simple_mapping)
            log(f"{key}: created slide {new_index+1} tokens after dup: {tokens_in_slide(new_slide)}")
            created_indices.append(new_index)
            current_index = new_index

        # Now we have [seed_index] + created_indices as our sequence in order
        sequence_indices = [seed_index] + created_indices
        log(f"{key}: sequence slide indices: {sequence_indices}")

        # Replace the body token with each chunk on seed + duplicates
        for seq_i, chunk_text in zip(sequence_indices, chunks):
            slide = prs.slides[seq_i]
            replace_tokens_in_slide(slide, {key: chunk_text})
            preview = (chunk_text or "")[:80].replace('\n','|')
            log(f"{key}: slide {seq_i+1} text set preview: {preview}...")

    # Snapshot after all replacements and duplications
    snapshot('after')

    # Write output
    # Optionally stamp output filename with current date-time
    out_path_str = args.out
    if args.stamp:
        base, ext = os.path.splitext(out_path_str)
        ts = datetime.now().strftime('%Y%m%d-%H%M%S')
        out_path_str = f"{base}.{ts}{ext or '.pptx'}"

    # Update core properties for traceability
    try:
        now = datetime.now()
        prs.core_properties.modified = now
        prs.core_properties.last_modified_by = "auto-lectio"
    except Exception:
        pass

    out_path = Path(out_path_str)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    # Always print the final output path so users can find the newest
    print(f"Wrote: {out_path}")


if __name__ == "__main__":
    main()
