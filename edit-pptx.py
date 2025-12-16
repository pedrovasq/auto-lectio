from pptx import Presentation

prs = Presentation('practice.pptx')

START_TOKEN = '{{FIRST_READING}}'

def iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        
        if shape.shape_type == 6:
            for s in shape.shapes:
                yield s

def find_token(prs, token):
    for si, slide in enumerate(prs.slides):
        for shape in iter_shapes(slide):
            if shape.has_text_frame and token in shape.text_frame.text:
                return si, slide, shape
    return None, None, None

def replace_token_in_shape(shape, token, new_text):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            if token in run.text:
                run.text = run.text.replace(token, new_text)

start_idx, start_slide, start_shape = find_token(prs, START_TOKEN)
if start_slide is None:
    raise RuntimeError(f"Could not find {START_TOKEN}")

replace_token_in_shape(start_shape, START_TOKEN, "THIS WORKED! BUT WILL IT WRAP PROPERLY WHO KNOWS CUZ I DONT!")

prs.save('practice-out.pptx')
